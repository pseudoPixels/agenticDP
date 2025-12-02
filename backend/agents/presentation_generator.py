from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
import json
import re
import os
import random
import io
from typing import Dict, List, Any, Optional
from .image_generator import ImageGeneratorAgent


class PresentationGeneratorAgent:
    """Agent responsible for generating professional presentation decks with images"""
    
    def __init__(self, api_key: str):
        self.client = genai.Client(api_key=api_key)
        self.model_name = 'gemini-2.0-flash-exp'
        self.image_generator = ImageGeneratorAgent(api_key)
        
        # Get available templates
        self.template_dir = os.path.join(os.path.dirname(__file__), 'slideTemplates')
        self.templates = self._get_available_templates()
        
    def _get_available_templates(self) -> List[str]:
        """Get list of available PPTX templates"""
        if not os.path.exists(self.template_dir):
            return []
        return [f for f in os.listdir(self.template_dir) if f.endswith('.pptx')]
    
    def _select_random_template(self) -> Optional[str]:
        """Select a random template from available templates"""
        if not self.templates:
            return None
        template_name = random.choice(self.templates)
        return os.path.join(self.template_dir, template_name)
    
    def generate_presentation(self, topic: str) -> Dict[str, Any]:
        """Generate a comprehensive presentation structure on the given topic"""
        
        prompt = f"""You are an expert presentation designer. Create a professional, engaging presentation on: "{topic}"

Structure the presentation with the following slides (return as valid JSON):
{{
    "title": "Presentation title",
    "subtitle": "Brief engaging subtitle",
    "slides": [
        {{
            "type": "title",
            "title": "Main title",
            "content": "Subtitle or tagline",
            "image_prompt": "Detailed description for an educational/professional image"
        }},
        {{
            "type": "section",
            "title": "Section heading",
            "image_prompt": "Detailed description for section image"
        }},
        {{
            "type": "content",
            "title": "Slide title",
            "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
            "image_prompt": "Detailed description for content image"
        }},
        {{
            "type": "chart",
            "title": "Chart slide title",
            "chart_data": {{
                "type": "bar",
                "title": "Chart title",
                "categories": ["Category 1", "Category 2", "Category 3"],
                "values": [10, 20, 30]
            }},
            "image_prompt": "Detailed description for supporting image"
        }},
        {{
            "type": "closing",
            "title": "Thank You!",
            "content": "Contact information or closing message",
            "image_prompt": "Celebratory or thank you themed image"
        }}
    ]
}}

IMPORTANT GUIDELINES:
- Create 8-12 slides total
- Start with a title slide
- Include 1-2 section divider slides
- Include 5-7 content slides with bullet points (3-5 bullets each)
- Optionally include 1 chart slide if data visualization would help
- End with a closing/thank you slide
- EVERY slide must have an image_prompt for visual appeal
- Make content professional, clear, and engaging
- Keep bullet points concise (1-2 lines each)

Return ONLY the JSON, no markdown formatting or extra text."""

        try:
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=[prompt]
            )
            
            # Extract text from response
            presentation_text = ""
            if hasattr(response, 'text') and response.text:
                presentation_text = response.text
            elif hasattr(response, 'candidates') and response.candidates:
                for candidate in response.candidates:
                    if hasattr(candidate, 'content') and candidate.content:
                        for part in candidate.content.parts:
                            if hasattr(part, 'text') and part.text:
                                presentation_text += part.text
            
            presentation_text = presentation_text.strip()
            
            # Clean up markdown formatting if present
            presentation_text = re.sub(r'^```json\s*', '', presentation_text)
            presentation_text = re.sub(r'\s*```$', '', presentation_text)
            
            presentation_data = json.loads(presentation_text)
            
            # Add metadata
            presentation_data['topic'] = topic
            presentation_data['version'] = 1
            
            return presentation_data
            
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {e}")
            print(f"Response text: {presentation_text}")
            return self._create_fallback_presentation(topic)
        except Exception as e:
            print(f"Error generating presentation: {e}")
            import traceback
            traceback.print_exc()
            return self._create_fallback_presentation(topic)
    
    def _create_fallback_presentation(self, topic: str) -> Dict[str, Any]:
        """Create a basic presentation structure if generation fails"""
        return {
            "title": f"Presentation: {topic}",
            "subtitle": "A comprehensive overview",
            "topic": topic,
            "version": 1,
            "slides": [
                {
                    "type": "title",
                    "title": topic,
                    "content": "A comprehensive overview",
                    "image_prompt": f"Professional image about {topic}"
                },
                {
                    "type": "content",
                    "title": "Overview",
                    "content": [
                        f"Introduction to {topic}",
                        "Key concepts and ideas",
                        "Practical applications"
                    ],
                    "image_prompt": f"Educational illustration about {topic}"
                },
                {
                    "type": "closing",
                    "title": "Thank You!",
                    "image_prompt": "Thank you celebration image"
                }
            ]
        }
    
    def create_pptx(self, presentation_data: Dict[str, Any], images: Dict[str, str]) -> io.BytesIO:
        """Create a PPTX file from presentation data and images - matching HTML exactly"""
        
        # Create a blank presentation (no template)
        prs = Presentation()
        prs.slide_width = Inches(10)  # 16:9 aspect ratio
        prs.slide_height = Inches(5.625)
        # Define colors to match HTML
        primary_color = RGBColor(16, 185, 129)  # emerald-500
        accent_color = RGBColor(20, 184, 166)   # teal-500
        text_color = RGBColor(17, 24, 39)       # gray-900
        
        # Create slides
        slides_data = presentation_data.get('slides', [])
        for i, slide_info in enumerate(slides_data):
            slide_type = slide_info.get('type', 'content')
            image_key = f"slide_{i}"
            image_data = images.get(image_key)
            
            print(f"   Creating slide {i+1}/{len(slides_data)}: {slide_type}")
            
            if slide_type == 'title':
                self._create_title_slide_with_image(prs, slide_info, image_data, primary_color, accent_color)
            elif slide_type == 'section':
                self._create_section_slide_with_image(prs, slide_info, image_data, primary_color, accent_color)
            elif slide_type == 'content':
                self._create_content_slide_with_image(prs, slide_info, image_data, primary_color, text_color, accent_color)
            elif slide_type == 'chart':
                self._create_chart_slide_with_image(prs, slide_info, image_data, primary_color, accent_color)
            elif slide_type == 'closing':
                self._create_closing_slide_with_image(prs, slide_info, image_data, primary_color, accent_color)
            else:
                # Default to content slide
                self._create_content_slide_with_image(prs, slide_info, image_data, primary_color, text_color, accent_color)
        
        # Save to BytesIO
        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        
        return pptx_stream
    
    def _generate_image(self, prompt: str) -> Optional[io.BytesIO]:
        """Generate image from prompt and return as BytesIO"""
        try:
            base64_image = self.image_generator.generate_image(prompt, style='realistic')
            if base64_image:
                # Convert base64 to BytesIO
                import base64
                if ',' in base64_image:
                    base64_image = base64_image.split(',')[1]
                image_bytes = base64.b64decode(base64_image)
                return io.BytesIO(image_bytes)
        except Exception as e:
            print(f"Error generating image: {e}")
        return None
    
    # ---------- Slide Builders with Images (Using Template Layouts) ----------
    
    def _create_title_slide_with_image(self, prs, slide_info, image_data, primary_color, accent_color):
        """Create title slide matching HTML - emerald to teal gradient"""
        # Create blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Add gradient background (emerald to teal)
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = RGBColor(16, 185, 129)  # emerald-500
        fill.gradient_stops[1].color.rgb = RGBColor(20, 184, 166)  # teal-500

        title_text = slide_info.get("title", "Title")
        content = slide_info.get("content", "")
        if isinstance(content, list):
            content = " | ".join(str(x) for x in content)

        # Add title - centered at top
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Add subtitle if exists
        if content:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.6))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = str(content)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Add image at bottom if available
        if image_data:
            image_stream = self._base64_to_stream(image_data)
            if image_stream:
                left, top, width = Inches(2.5), Inches(3.5), Inches(5)
                slide.shapes.add_picture(image_stream, left, top, width=width)
                print(f"      ✅ Image added to title slide")

    def _create_section_slide_with_image(self, prs, slide_info, image_data, primary_color, accent_color):
        """Create section slide matching HTML - blue to indigo gradient"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add gradient background (blue to indigo)
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = RGBColor(59, 130, 246)  # blue-500
        fill.gradient_stops[1].color.rgb = RGBColor(99, 102, 241)  # indigo-500

        title_text = slide_info.get("title", "Section")

        # Add title - centered
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Add image if available
        if image_data:
            image_stream = self._base64_to_stream(image_data)
            if image_stream:
                left, top, width = Inches(2.5), Inches(3), Inches(5)
                slide.shapes.add_picture(image_stream, left, top, width=width)
                print(f"      ✅ Image added to section slide")

    def _create_content_slide_with_image(self, prs, slide_info, image_data, primary_color, text_color, accent_color):
        """Create content slide matching HTML - white background with border"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title_text = slide_info.get("title", "Slide")
        content_items = slide_info.get("content", [])
        if not isinstance(content_items, list):
            content_items = [str(content_items)]

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)  # gray-900

        # Add bullet points on left
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = f"• {str(item)}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(55, 65, 81)  # gray-700
            p.space_before = Pt(12)

        # Add image on right if available
        if image_data:
            image_stream = self._base64_to_stream(image_data)
            if image_stream:
                left, top, width = Inches(5.5), Inches(1.5), Inches(4)
                slide.shapes.add_picture(image_stream, left, top, width=width)
                print(f"      ✅ Image added to content slide")

    def _create_chart_slide_with_image(self, prs, slide_info, image_data, primary_color, accent_color):
        """Create chart slide matching HTML - white background"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title_text = slide_info.get("title", "Chart")

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(17, 24, 39)

        # Add chart on left side
        chart_data_info = slide_info.get("chart_data", {}) or {}
        chart_type_str = (chart_data_info.get("type") or "bar").lower()
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()
        cats = chart_data_info.get("categories") or ["A", "B", "C"]
        vals = chart_data_info.get("values") or [10, 20, 30]
        chart_data.categories = cats
        chart_data.add_series("Data", vals)

        # Chart positioned on left side
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(4.5), Inches(4)
        chart_shape = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = chart_shape.chart

        title_txt = chart_data_info.get("title")
        chart.has_title = bool(title_txt)
        if title_txt:
            chart.chart_title.text_frame.text = str(title_txt)
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)

        chart.has_legend = True

        # Add image on right if available
        if image_data:
            image_stream = self._base64_to_stream(image_data)
            if image_stream:
                left, top, width = Inches(5.5), Inches(1.5), Inches(4)
                slide.shapes.add_picture(image_stream, left, top, width=width)
                print(f"      ✅ Image added to chart slide")

    def _create_closing_slide_with_image(self, prs, slide_info, image_data, primary_color, accent_color):
        """Create closing slide matching HTML - purple to pink gradient"""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add gradient background (purple to pink)
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = RGBColor(168, 85, 247)  # purple-500
        fill.gradient_stops[1].color.rgb = RGBColor(236, 72, 153)  # pink-500

        title_text = slide_info.get("title", "Thank You!")
        content = slide_info.get("content")
        if isinstance(content, list):
            content = "\n".join(str(x) for x in content)

        # Add title - centered
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Add content if exists
        if content:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.6))
            content_frame = content_box.text_frame
            content_frame.text = str(content)
            content_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            content_frame.paragraphs[0].font.size = Pt(24)
            content_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Add image if available
        if image_data:
            image_stream = self._base64_to_stream(image_data)
            if image_stream:
                left, top, width = Inches(2.5), Inches(3.5), Inches(5)
                slide.shapes.add_picture(image_stream, left, top, width=width)
                print(f"      ✅ Image added to closing slide")
    
    def _base64_to_stream(self, base64_data: str) -> Optional[io.BytesIO]:
        """Convert base64 image data to BytesIO stream"""
        try:
            import base64
            if ',' in base64_data:
                base64_data = base64_data.split(',')[1]
            image_bytes = base64.b64decode(base64_data)
            return io.BytesIO(image_bytes)
        except Exception as e:
            print(f"Error converting base64 to stream: {e}")
            return None
